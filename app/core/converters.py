import fitz  # PyMuPDF
from PIL import Image
import os
import io
from pypdf import PdfWriter, PdfReader
import tempfile
from urllib.parse import urlparse
from pathlib import Path
from app.core.logger import get_logger

LOG = get_logger(__name__)

import sys

# NOTA: las dependencias pesadas (pandas, pdf2docx, python-pptx, comtypes,
# docx2pdf, pytesseract, requests, img2pdf) se importan de forma diferida
# DENTRO de las funciones que las usan, para acelerar el arranque de la app
# (que importa este módulo al cargar varias páginas).

_TESSERACT_CONFIGURED = False


def _ensure_tesseract():
    """Configura la ruta de Tesseract la primera vez que se necesita OCR."""
    global _TESSERACT_CONFIGURED
    import pytesseract
    if not _TESSERACT_CONFIGURED:
        pytesseract.pytesseract.tesseract_cmd = get_tesseract_cmd()
        _TESSERACT_CONFIGURED = True
    return pytesseract


# Configure Tesseract path
def get_tesseract_cmd():
    if getattr(sys, 'frozen', False):
        # Search for tesseract in internal temporary folder (_MEIPASS)
        bundled_path = os.path.join(sys._MEIPASS, "Tesseract-OCR", "tesseract.exe")
        if os.path.exists(bundled_path):
            return bundled_path
        
        # Or search next to the executable
        local_path = os.path.join(os.path.dirname(sys.executable), "Tesseract-OCR", "tesseract.exe")
        if os.path.exists(local_path):
            return local_path
            
    # In Development: Check if there is a local Tesseract-OCR folder in the project
    dev_local_path = os.path.join(os.getcwd(), "Tesseract-OCR", "tesseract.exe")
    if os.path.exists(dev_local_path):
        return dev_local_path
        
    # In bin/Tesseract-OCR folder (Portable)
    bin_sub_path = os.path.join(os.getcwd(), "bin", "Tesseract-OCR", "tesseract.exe")
    if os.path.exists(bin_sub_path):
        return bin_sub_path
        
    # Directly in bin (Portable)
    bin_root_path = os.path.join(os.getcwd(), "bin", "tesseract.exe")
    if os.path.exists(bin_root_path):
        return bin_root_path

    # By default: Assume it is in the system PATH
    return "tesseract"

class Converter:
    @staticmethod
    def images_to_pdf(image_list, output_path):
        """
        Convert a list of images into a single PDF.
        :param image_list: List of paths to image files
        :param output_path: Path to save the generated PDF
        """
        import img2pdf
        with open(output_path, "wb") as f:
            f.write(img2pdf.convert(image_list))

    @staticmethod
    def pdf_to_images(pdf_path, output_dir, fmt='png', pages=None):
        """
        Convert a PDF into a series of images.
        :param pdf_path: Path to the PDF file
        :param output_dir: Directory to save the images
        :param fmt: Image format (default: png)
        :param pages: List of page indices (starting with 0) to convert. If None, converts all.
        """
        if not os.path.exists(output_dir):
            os.makedirs(output_dir)
            
        doc = fitz.open(pdf_path)
        base_name = os.path.splitext(os.path.basename(pdf_path))[0]
        
        for i, page in enumerate(doc):
            if pages is not None and i not in pages:
                continue
                
            pix = page.get_pixmap()
            output_filename = f"{base_name}_page_{i+1}.{fmt}"
            output_path = os.path.join(output_dir, output_filename)
            pix.save(output_path)
        
        doc.close()

    @staticmethod
    def pdf_to_word(pdf_path, output_path):
        """
        Convert PDF to Word (DOCX).
        :param pdf_path: Path to the PDF file
        :param output_path: Path to save the DOCX file
        """
        from pdf2docx import Converter as DocxConverter
        cv = DocxConverter(pdf_path)
        cv.convert(output_path, start=0, end=None)
        cv.close()

    @staticmethod
    def pdf_to_excel(pdf_path, output_path):
        """
        Convert PDF to Excel (XLSX). Extracts tables from each page and writes
        them to separate sheets. Uses the first row as header when it looks
        like one (all cells non-empty and the row below has data).
        :param pdf_path: Path to the PDF file
        :param output_path: Path to save the XLSX file
        """
        import pdfplumber
        import pandas as pd

        def _looks_like_header(rows):
            if len(rows) < 2:
                return False
            head = rows[0]
            # Cabecera plausible: todas las celdas con texto y sin valores vacíos.
            return all(c is not None and str(c).strip() != "" for c in head)

        with pdfplumber.open(pdf_path) as pdf:
            with pd.ExcelWriter(output_path, engine='openpyxl') as writer:
                has_tables = False
                for i, page in enumerate(pdf.pages):
                    tables = page.extract_tables()
                    if not tables:
                        continue
                    has_tables = True
                    for j, table in enumerate(tables):
                        if _looks_like_header(table):
                            df = pd.DataFrame(table[1:], columns=table[0])
                            use_header = True
                        else:
                            df = pd.DataFrame(table)
                            use_header = False
                        # Nombre de hoja válido (máx 31 chars, sin caracteres prohibidos).
                        sheet_name = f"Pag{i+1}_Tabla{j+1}"[:31]
                        df.to_excel(writer, sheet_name=sheet_name, index=False, header=use_header)

                if not has_tables:
                    pd.DataFrame(["No se encontraron tablas"]).to_excel(
                        writer, sheet_name="Info", index=False, header=False)

    @staticmethod
    def pdf_to_powerpoint(pdf_path, output_path):
        """
        (Function not implemented in the GUI)
        Convert PDF to PowerPoint (PPTX).
        Converts pages into images and places them on slides.
        :param pdf_path: Path to the PDF file
        :param output_path: Path to save the PPTX file
        """
        from pptx import Presentation
        prs = Presentation()
        # Use a temporary directory for images
        
        with tempfile.TemporaryDirectory() as temp_dir:
            doc = fitz.open(pdf_path)
            
            for i, page in enumerate(doc):
                # Render page to image
                pix = page.get_pixmap()
                image_path = os.path.join(temp_dir, f"pagina_{i}.png")
                pix.save(image_path)
                
                # Add a blank slide
                blank_slide_layout = prs.slide_layouts[6] 
                slide = prs.slides.add_slide(blank_slide_layout)
                
                # Update presentation slide size based on dimensions of the first page
                if i == 0:
                    # page.rect gives size in points (1/72 inch)
                    width_pt = page.rect.width
                    height_pt = page.rect.height
                    
                    # Convert points to EMU (1 point = 12700 EMU)
                    prs.slide_width = int(width_pt * 12700)
                    prs.slide_height = int(height_pt * 12700)
                
                # Add image to slide
                slide.shapes.add_picture(image_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
            
            doc.close()
        prs.save(output_path)

    @staticmethod
    def _find_soffice():
        """Localiza el ejecutable de LibreOffice (soffice) si está instalado."""
        import shutil
        for name in ("soffice", "soffice.exe"):
            found = shutil.which(name)
            if found:
                return found
        candidates = [
            r"C:\Program Files\LibreOffice\program\soffice.exe",
            r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
            "/usr/bin/soffice",
            "/usr/bin/libreoffice",
            "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        ]
        for path in candidates:
            if os.path.exists(path):
                return path
        return None

    @staticmethod
    def _libreoffice_to_pdf(input_path, output_path):
        """Convierte cualquier documento de Office a PDF usando LibreOffice
        en modo headless. Alternativa multiplataforma cuando MS Office no está
        disponible.
        """
        import subprocess
        soffice = Converter._find_soffice()
        if not soffice:
            raise RuntimeError(
                "No se encontró Microsoft Office ni LibreOffice instalados.\n"
                "Instala alguno de ellos para convertir documentos de Office a PDF."
            )

        out_dir = tempfile.mkdtemp(prefix="pf_office_")
        try:
            subprocess.run(
                [soffice, "--headless", "--convert-to", "pdf", "--outdir", out_dir, os.path.abspath(input_path)],
                check=True, timeout=180,
                stdout=subprocess.PIPE, stderr=subprocess.PIPE,
            )
            produced = os.path.join(
                out_dir, os.path.splitext(os.path.basename(input_path))[0] + ".pdf"
            )
            if not os.path.exists(produced):
                raise RuntimeError("LibreOffice no generó el PDF esperado.")
            os.replace(produced, output_path)
        finally:
            import shutil
            shutil.rmtree(out_dir, ignore_errors=True)

    @staticmethod
    def word_to_pdf(doc_path, output_path):
        """
        Convert Word to PDF. Tries MS Word (docx2pdf) first and falls back to
        LibreOffice if Word is not installed.
        :param doc_path: Path to DOCX/DOC file
        :param output_path: Path to save the PDF file
        """
        try:
            from docx2pdf import convert as docx_convert
            docx_convert(doc_path, output_path)
            if os.path.exists(output_path):
                return
            raise RuntimeError("docx2pdf no produjo ningún archivo.")
        except Exception as e:
            LOG.warning(f"Word->PDF con MS Office falló ({e}); intentando LibreOffice...")
            Converter._libreoffice_to_pdf(doc_path, output_path)

    @staticmethod
    def excel_to_pdf(excel_path, output_path):
        """
        Convert Excel to PDF. Tries MS Excel (COM) first and falls back to
        LibreOffice if Excel is not installed.
        :param excel_path: Path to XLSX/XLS file
        :param output_path: Path to save the PDF file
        """
        import comtypes.client
        excel_path = os.path.abspath(excel_path)
        output_path = os.path.abspath(output_path)

        try:
            excel = comtypes.client.CreateObject("Excel.Application")
            excel.Visible = False
            try:
                wb = excel.Workbooks.Open(excel_path)
                # 0 is xlTypePDF
                wb.ExportAsFixedFormat(0, output_path)
                wb.Close()
            finally:
                excel.Quit()
            if os.path.exists(output_path):
                return
            raise RuntimeError("Excel no produjo ningún archivo.")
        except Exception as e:
            LOG.warning(f"Excel->PDF con MS Office falló ({e}); intentando LibreOffice...")
            Converter._libreoffice_to_pdf(excel_path, output_path)

    @staticmethod
    def powerpoint_to_pdf(ppt_path, output_path):
        """
        Convert PowerPoint to PDF. Tries MS PowerPoint (COM) first and falls
        back to LibreOffice if PowerPoint is not installed.
        :param ppt_path: Path to PPTX/PPT file
        :param output_path: Path to save the PDF file
        """
        import comtypes.client
        ppt_path = os.path.abspath(ppt_path)
        output_path = os.path.abspath(output_path)

        try:
            powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
            try:
                deck = powerpoint.Presentations.Open(ppt_path)
                # 32 is ppSaveAsPDF
                deck.SaveAs(output_path, 32)
                deck.Close()
            finally:
                powerpoint.Quit()
            if os.path.exists(output_path):
                return
            raise RuntimeError("PowerPoint no produjo ningún archivo.")
        except Exception as e:
            LOG.warning(f"PPT->PDF con MS Office falló ({e}); intentando LibreOffice...")
            Converter._libreoffice_to_pdf(ppt_path, output_path)

    @staticmethod
    def apply_ocr(pdf_path, output_path, lang='eng+spa'):
        """
        Apply OCR to a scanned PDF to make it searchable.
        :param pdf_path: Path to scanned PDF
        :param output_path: Path to save searchable PDF
        :param lang: Language code for Tesseract (default: 'eng+spa')
        """
        # check if tesseract is installed/available
        pytesseract = _ensure_tesseract()

        doc = fitz.open(pdf_path)
        writer = PdfWriter()

        for i, page in enumerate(doc):
            # Get page image
            # matrix=fitz.Matrix(2, 2) for better resolution/precision
            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
            img_bytes = pix.tobytes("png")
            img = Image.open(io.BytesIO(img_bytes))
            
            # Get PDF data from tesseract
            # this returns a PDF file as bytes
            try:
                pdf_bytes = pytesseract.image_to_pdf_or_hocr(img, extension='pdf', lang=lang)
            except pytesseract.TesseractNotFoundError:
                raise Exception("Tesseract no está instalado o no está en su PATH. Por favor instale Tesseract-OCR.")
                
            # Create a PDF reader from these bytes
            ocr_page_reader = PdfReader(io.BytesIO(pdf_bytes))
            writer.add_page(ocr_page_reader.pages[0])
            
        doc.close()
        
        with open(output_path, "wb") as f:
            writer.write(f)

    @staticmethod
    def _is_url(s: str) -> bool:
        try:
            p = urlparse(s)
            return p.scheme in ("http", "https")
        except Exception:
            return False

    @staticmethod
    def _is_html_string(s: str) -> bool:
        return s.strip().lower().startswith("<!doctype") or s.strip().lower().startswith("<html")

    @staticmethod
    def html_to_pdf_playwright(source: str, output_path: str, *, is_url: bool = None,
                               paper_format: str = "A4", landscape: bool = False,
                               margin_mm: int = 12, scale: float = 1.0,
                               wait_for_network_idle: bool = True, timeout: int = 60_000):
        """
        Converts HTML -> PDF with high fidelity Playwright (headless Chromium).
        """
        output_path = str(output_path)
        # Determine source type if not specified
        if is_url is None:
            is_url = Converter._is_url(source)

        page_url = ""
        temp_file = None

        if not is_url and os.path.exists(source) and Path(source).suffix.lower() in (".html", ".htm"):
            page_url = Path(source).absolute().as_uri()
        elif not is_url and Converter._is_html_string(source):
            temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=".html", mode="w", encoding="utf-8")
            temp_file.write(source)
            temp_file.flush()
            temp_file.close()
            page_url = Path(temp_file.name).absolute().as_uri()
        elif is_url:
            page_url = source
        else:
            if os.path.exists(source):
                page_url = Path(source).absolute().as_uri()
            else:
                raise ValueError("Error en la conversion HTML->PDF. Revisar rutas o archivos.")

        # Playwright Check
        try:
            from playwright.sync_api import sync_playwright, TimeoutError as PWTimeout
        except ImportError as e:
            raise ImportError("Playwright no instalado. Ejecuta 'pip install playwright' y 'playwright install'.") from e

        mm_to_in = lambda mm: mm / 25.4

        try:
            with sync_playwright() as p:
                browser = p.chromium.launch(headless=True, args=["--no-sandbox", "--disable-setuid-sandbox"])
                context = browser.new_context(viewport={"width": 1280, "height": 800})
                page = context.new_page()
                
                nav_opts = {"timeout": timeout}
                nav_opts["wait_until"] = "networkidle" if wait_for_network_idle else "load"
                
                page.goto(page_url, **nav_opts)
                
                # Wait for canvas/graphics
                try:
                    page.wait_for_function(
                        """
                        () => {
                            const canvases = Array.from(document.querySelectorAll('canvas'));
                            if (canvases.length === 0) return true;
                            return canvases.every(c => {
                                try { return c.toDataURL().length > 200; } catch(e) { return false; }
                            });
                        }
                        """,
                        timeout=5000
                    )
                except Exception:
                    pass
                
                page.wait_for_timeout(500)

                pdf_options = {
                    "path": output_path,
                    "format": paper_format,
                    "print_background": True,
                    "landscape": landscape,
                    "scale": scale,
                    "margin": {
                        "top": f"{mm_to_in(margin_mm)}in",
                        "bottom": f"{mm_to_in(margin_mm)}in",
                        "left": f"{mm_to_in(margin_mm)}in",
                        "right": f"{mm_to_in(margin_mm)}in",
                    }
                }
                
                page.pdf(**pdf_options)
                context.close()
                browser.close()
                
        finally:
            if temp_file and os.path.exists(temp_file.name):
                try:
                    os.unlink(temp_file.name)
                except Exception:
                    pass

    @staticmethod
    def pdf_to_text(pdf_path, output_path, preserve_layout=True):
        """
        Extract text from PDF to a TXT file.
        :param pdf_path: Path to the PDF file
        :param output_path: Path to save the TXT file
        :param preserve_layout: If True, tries to maintain physical layout (fitz 'text' or 'blocks'). 
                                If False, only raw string (may lose line breaks).
        """
        doc = fitz.open(pdf_path)
        full_text = []

        for page in doc:
            # "text" preserves lines and paragraphs approximately
            text = page.get_text("text") if preserve_layout else page.get_text("text").replace("\n", " ")
            full_text.append(text)
        
        doc.close()

        with open(output_path, "w", encoding="utf-8") as f:
            f.write("\n\n".join(full_text))

    @staticmethod
    def html_to_pdf(source, output_path, is_url=False):
        """
        Convert HTML (File or URL) to PDF.
        :param source: URL string or file path
        :param output_path: Output PDF path
        :param is_url: Boolean
        """
        pdf_bytes = None
        
        if is_url:
            try:
                import requests
                response = requests.get(source)
                response.raise_for_status()
                html_data = response.content
                # fitz.open("pdf", data) expects pdf data.
                # PyMuPDF v1.24+ supports "story" but standard fitz.open(stream=html) is for opening, 
                # then convert_to_pdf().
                
                # Check if we can open html directly from memory
                doc = fitz.open("html", html_data)
                pdf_bytes = doc.convert_to_pdf()
                doc.close()
            except Exception as e:
                raise Exception(f"Error al obtener o convertir la URL: {e}")
        else:
            # Local file
            if not os.path.exists(source):
                raise FileNotFoundError(f"Archivo HTML no encontrado: {source}")
            
            doc = fitz.open(source)
            pdf_bytes = doc.convert_to_pdf()
            doc.close()
            
        if pdf_bytes:
            with open(output_path, "wb") as f:
                f.write(pdf_bytes)

    @staticmethod
    def pdf_to_html(pdf_path, output_path, mode="text", progress_callback=None, cancel_check=None):
        """
        Convert PDF to HTML.
        :param pdf_path: Input PDF
        :param output_path: Output HTML
        :param mode: "text" (Pres. layout), "visual" (Images), "ocr" (Tesseract HOCR)
        :param progress_callback: Function(current_page, total_pages)
        :param cancel_check: Function() -> bool. If returns True, cancels.
        """
        # Compatibility wrapper for old calls if any
        if isinstance(mode, bool):
             mode = "ocr" if mode else "text"

        if mode == "text":
            # Simple Mode: Extract text/images as HTML
            doc = fitz.open(pdf_path)
            total_pages = len(doc)
            
            # We will generate a single HTML with all pages added.
            html_content = ["<html><body>"]
            for i, page in enumerate(doc):
                if cancel_check and cancel_check():
                    doc.close()
                    return # Cancelled
                
                html_content.append(page.get_text("html"))
                html_content.append("<hr>") # Page break separator
                
                if progress_callback:
                    progress_callback(i + 1, total_pages)
            
            html_content.append("</body></html>")
            doc.close()
            
            with open(output_path, "w", encoding="utf-8") as f:
                f.write("\n".join(html_content))
                
        elif mode == "visual":
            # Visual Mode: Render pages as images (Base64) to preserve perfect look
            doc = fitz.open(pdf_path)
            total_pages = len(doc)
            html_content = ["<html><head><style>body { background: #525659; margin: 0; padding: 20px; text-align: center; } img { box-shadow: 0 0 10px rgba(0,0,0,0.5); margin-bottom: 20px; max-width: 100%; }</style></head><body>"]
            
            import base64
            
            for i, page in enumerate(doc):
                if cancel_check and cancel_check():
                    doc.close()
                    return
                
                # Render high quality (zoom 2x)
                pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                img_data = pix.tobytes("png")
                b64_img = base64.b64encode(img_data).decode('utf-8')
                
                html_content.append(f'<div class="page"><img src="data:image/png;base64,{b64_img}" alt="Page {i+1}" /></div>')
                
                if progress_callback:
                    progress_callback(i + 1, total_pages)
            
            html_content.append("</body></html>")
            doc.close()
            
            with open(output_path, "w", encoding="utf-8") as f:
                f.write("\n".join(html_content))

        elif mode == "ocr":
            # Complex Mode
            pytesseract = _ensure_tesseract()
            doc = fitz.open(pdf_path)
            total_pages = len(doc)
            full_hocr = ["<html><body>"]
            
            for i, page in enumerate(doc):
                if cancel_check and cancel_check():
                    doc.close()
                    return
                    
                pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                img_bytes = pix.tobytes("png")
                pil_img = Image.open(io.BytesIO(img_bytes))
                
                try:
                    hocr = pytesseract.image_to_pdf_or_hocr(pil_img, extension='hocr', lang='eng+spa')
                    full_hocr.append(hocr.decode('utf-8'))
                except Exception:
                    full_hocr.append(f"<p>Error de OCR en la página {i+1}</p>")
                    
                full_hocr.append("<hr>")
                
                if progress_callback:
                    progress_callback(i + 1, total_pages)
            
            full_hocr.append("</body></html>")
            doc.close()
            
            with open(output_path, "w", encoding="utf-8") as f:
                f.write("\n".join(full_hocr))
